"""Concrete loader implementations for every supported file type."""
import csv
import io
import json
import logging
from pathlib import Path

from app.loaders.base import (
    BaseDocumentLoader,
    DocumentLoadError,
    LoadedSegment,
    decode_text,
)

logger = logging.getLogger(__name__)

# Extensions treated as plain text / source code, mapped to a language label.
CODE_EXTENSIONS: dict[str, str] = {
    ".py": "python", ".js": "javascript", ".jsx": "javascript", ".ts": "typescript",
    ".tsx": "typescript", ".java": "java", ".cpp": "cpp", ".cc": "cpp", ".cxx": "cpp",
    ".c": "c", ".h": "c", ".hpp": "cpp", ".go": "go", ".rs": "rust", ".php": "php",
    ".sql": "sql", ".rb": "ruby", ".swift": "swift", ".kt": "kotlin", ".cs": "csharp",
    ".sh": "bash", ".ps1": "powershell", ".r": "r", ".scala": "scala", ".lua": "lua",
    ".yaml": "yaml", ".yml": "yaml", ".toml": "toml", ".ini": "ini", ".cfg": "ini",
}


class PdfLoader(BaseDocumentLoader):
    extensions = (".pdf",)

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            segments = []
            for i, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    segments.append(LoadedSegment(content=text, page=i))
            return segments
        except DocumentLoadError:
            raise
        except Exception as exc:
            raise DocumentLoadError(f"Failed to parse PDF: {exc}") from exc


class TextLoader(BaseDocumentLoader):
    extensions = (".txt", ".log")

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        text = decode_text(self._read_bytes(file_path)).strip()
        return [LoadedSegment(content=text)] if text else []


class MarkdownLoader(BaseDocumentLoader):
    extensions = (".md", ".markdown", ".mdx")

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        text = decode_text(self._read_bytes(file_path))
        # Split on top-level / second-level headings to keep section metadata.
        segments: list[LoadedSegment] = []
        current_title: str | None = None
        current_lines: list[str] = []

        def flush() -> None:
            body = "\n".join(current_lines).strip()
            if body:
                segments.append(LoadedSegment(content=body, section=current_title))

        for line in text.splitlines():
            if line.startswith(("# ", "## ")):
                flush()
                current_title = line.lstrip("#").strip()
                current_lines = [line]
            else:
                current_lines.append(line)
        flush()
        return segments or ([LoadedSegment(content=text.strip())] if text.strip() else [])


class CodeLoader(BaseDocumentLoader):
    extensions = tuple(CODE_EXTENSIONS.keys())

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        text = decode_text(self._read_bytes(file_path)).strip()
        if not text:
            return []
        language = CODE_EXTENSIONS.get(Path(file_path).suffix.lower(), "text")
        return [LoadedSegment(content=text, extra={"language": language})]


class CsvLoader(BaseDocumentLoader):
    extensions = (".csv", ".tsv")

    ROWS_PER_SEGMENT = 40

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        text = decode_text(self._read_bytes(file_path))
        delimiter = "\t" if str(file_path).lower().endswith(".tsv") else ","
        try:
            rows = list(csv.reader(io.StringIO(text), delimiter=delimiter))
        except csv.Error as exc:
            raise DocumentLoadError(f"Malformed CSV: {exc}") from exc
        if not rows:
            return []
        header = rows[0]
        segments = []
        for start in range(1, len(rows), self.ROWS_PER_SEGMENT):
            batch = rows[start : start + self.ROWS_PER_SEGMENT]
            lines = []
            for row_idx, row in enumerate(batch, start=start + 1):
                pairs = [f"{h}: {v}" for h, v in zip(header, row) if str(v).strip()]
                if pairs:
                    lines.append(f"Row {row_idx} — " + "; ".join(pairs))
            if lines:
                segments.append(
                    LoadedSegment(
                        content="\n".join(lines),
                        section=f"rows {start + 1}-{start + len(batch)}",
                    )
                )
        return segments


class ExcelLoader(BaseDocumentLoader):
    extensions = (".xlsx", ".xls", ".xlsm")

    ROWS_PER_SEGMENT = 40

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        try:
            import pandas as pd

            sheets = pd.read_excel(str(file_path), sheet_name=None, dtype=str)
        except Exception as exc:
            raise DocumentLoadError(f"Failed to parse Excel file: {exc}") from exc

        segments = []
        for sheet_name, df in sheets.items():
            df = df.fillna("")
            records = df.to_dict(orient="records")
            for start in range(0, len(records), self.ROWS_PER_SEGMENT):
                batch = records[start : start + self.ROWS_PER_SEGMENT]
                lines = []
                for row_idx, record in enumerate(batch, start=start + 2):  # +2: header row
                    pairs = [f"{k}: {v}" for k, v in record.items() if str(v).strip()]
                    if pairs:
                        lines.append(f"Row {row_idx} — " + "; ".join(pairs))
                if lines:
                    segments.append(
                        LoadedSegment(
                            content="\n".join(lines),
                            page=sheet_name,
                            section=f"Sheet {sheet_name}",
                        )
                    )
        return segments


class DocxLoader(BaseDocumentLoader):
    extensions = (".docx",)

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        try:
            import docx

            document = docx.Document(str(file_path))
        except Exception as exc:
            raise DocumentLoadError(f"Failed to parse Word document: {exc}") from exc

        segments: list[LoadedSegment] = []
        current_title: str | None = None
        current_parts: list[str] = []

        def flush() -> None:
            body = "\n".join(current_parts).strip()
            if body:
                segments.append(LoadedSegment(content=body, section=current_title))

        for para in document.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style is not None and (para.style.name or "").startswith("Heading"):
                flush()
                current_title = text
                current_parts = [text]
            else:
                current_parts.append(text)

        for table in document.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    current_parts.append(" | ".join(cells))
        flush()
        return segments


class PptxLoader(BaseDocumentLoader):
    extensions = (".pptx",)

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        try:
            from pptx import Presentation

            presentation = Presentation(str(file_path))
        except Exception as exc:
            raise DocumentLoadError(f"Failed to parse PowerPoint: {exc}") from exc

        segments = []
        for i, slide in enumerate(presentation.slides, start=1):
            texts = []
            title = None
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                shape_text = shape.text_frame.text.strip()
                if not shape_text:
                    continue
                if title is None and shape == getattr(slide.shapes, "title", None):
                    title = shape_text
                texts.append(shape_text)
            body = "\n".join(texts).strip()
            if body:
                segments.append(
                    LoadedSegment(content=body, page=i, section=title or f"Slide {i}")
                )
        return segments


class JsonLoader(BaseDocumentLoader):
    extensions = (".json", ".jsonl", ".ndjson")

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        text = decode_text(self._read_bytes(file_path)).strip()
        if not text:
            return []
        suffix = Path(file_path).suffix.lower()
        try:
            if suffix in (".jsonl", ".ndjson"):
                parsed = [json.loads(line) for line in text.splitlines() if line.strip()]
            else:
                parsed = json.loads(text)
        except json.JSONDecodeError as exc:
            raise DocumentLoadError(f"Malformed JSON: {exc}") from exc
        pretty = json.dumps(parsed, indent=2, ensure_ascii=False, default=str)
        return [LoadedSegment(content=pretty)]


class XmlLoader(BaseDocumentLoader):
    extensions = (".xml",)

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        raw = self._read_bytes(file_path)
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(raw, "xml")
            text = soup.get_text(separator="\n", strip=True)
        except Exception as exc:
            raise DocumentLoadError(f"Failed to parse XML: {exc}") from exc
        return [LoadedSegment(content=text)] if text else []


class HtmlLoader(BaseDocumentLoader):
    extensions = (".html", ".htm")

    def load(self, file_path: str | Path) -> list[LoadedSegment]:
        raw = self._read_bytes(file_path)
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(raw, "lxml")
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()
            title = soup.title.get_text(strip=True) if soup.title else None
            text = soup.get_text(separator="\n", strip=True)
        except Exception as exc:
            raise DocumentLoadError(f"Failed to parse HTML: {exc}") from exc
        return [LoadedSegment(content=text, section=title)] if text else []
